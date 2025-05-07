import streamlit as st
import os
import base64
from openai import OpenAI
from markitdown import MarkItDown
from pydantic import BaseModel, Field
from datetime import date
from enum import Enum
from typing import List
import io
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize OpenAI client for xAI API
try:
    client = OpenAI(
        api_key=os.getenv("XAI_API_KEY"),
        base_url="https://api.x.ai/v1",
    )
except Exception as e:
    st.error("Failed to initialize xAI API client. Please set the XAI_API_KEY environment variable.")
    logger.error(f"xAI API initialization error: {str(e)}")
    st.stop()

# Initialize MarkItDown
try:
    md = MarkItDown(enable_plugins=False)  # Plugins disabled by default
except ImportError as e:
    st.error("MarkItDown initialization failed. Ensure all dependencies are installed with: pip install 'markitdown[all]'")
    logger.error(f"MarkItDown initialization error: {str(e)}")
    st.stop()


    vendor_name: str = Field(description="Name of the vendor")
    vendor_address: Address = Field(description="Vendor's address")
    invoice_number: str = Field(description="Unique invoice identifier")
    invoice_date: date = Field(description="Date the invoice was issued")
    line_items: List[LineItem] = Field(description="List of purchased items/services")
    total_amount: float = Field(description="Total amount due", ge=0)
    currency: Currency = Field(description="Currency of the invoice")

class ContactInfo(BaseModel):
    name: str = Field(description="Full name of the contact")
    email: str = Field(description="Email address")
    phone: str = Field(description="Phone number")
    company: str = Field(description="Company name, if applicable")


# Helper to extract images from DOCX

def extract_images_from_docx(file):
    try:
        from docx import Document
        import tempfile
        import io
        images = []
        file.seek(0)
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
            tmp.write(file.read())
            tmp_path = tmp.name
        doc = Document(tmp_path)
        for rel in doc.part.rels.values():
            if 'image' in rel.target_ref:
                img_blob = rel.target_part.blob
                images.append(io.BytesIO(img_blob))
        return images
    except Exception as e:
        logger.error(f"DOCX image extraction error: {str(e)}")
        return []

# Helper to extract images from PPTX

def extract_images_from_pptx(file):
    try:
        from pptx import Presentation
        import tempfile
        import io
        images = []
        file.seek(0)
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            tmp.write(file.read())
            tmp_path = tmp.name
        prs = Presentation(tmp_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # PICTURE
                    image = shape.image
                    images.append(io.BytesIO(image.blob))
        return images
    except Exception as e:
        logger.error(f"PPTX image extraction error: {str(e)}")
        return []

# Helper to extract images from XLSX

def extract_images_from_xlsx(file):
    try:
        from openpyxl import load_workbook
        import tempfile
        import io
        images = []
        file.seek(0)
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
            tmp.write(file.read())
            tmp_path = tmp.name
        wb = load_workbook(tmp_path)
        for ws in wb.worksheets:
            for img in getattr(ws, '_images', []):
                if hasattr(img, 'ref') and hasattr(img, 'image'):
                    images.append(io.BytesIO(img.image))
        return images
    except Exception as e:
        logger.error(f"XLSX image extraction error: {str(e)}")
        return []

# Helper to convert Markdown to plain text

def markdown_to_text(md_content):
    try:
        import markdown2
        html = markdown2.markdown(md_content)
        import re
        text = re.sub('<[^<]+?>', '', html)
        return text
    except ImportError:
        import re
        text = re.sub(r'\*\*([^*]+)\*\*', r'\1', md_content)  # Remove bold
        text = re.sub(r'#+ ', '', text)  # Remove headers
        text = re.sub(r'!\[[^\]]*\]\([^)]*\)', '', text)  # Remove images
        text = re.sub(r'\[[^\]]*\]\([^)]*\)', '', text)  # Remove links
        text = re.sub(r'`([^`]*)`', r'\1', text)  # Remove inline code
        text = re.sub(r'\|.*\|', '', text)  # Remove table rows
        return text

# Helper to extract images from PDF

def extract_images_from_pdf(file):
    try:
        import fitz  # PyMuPDF
        import tempfile
        import io
        images = []
        file.seek(0)
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
            tmp.write(file.read())
            tmp_path = tmp.name
        doc = fitz.open(tmp_path)
        for page in doc:
            for img in page.get_images(full=True):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                images.append(io.BytesIO(image_bytes))
        return images
    except Exception as e:
        logger.error(f"PDF image extraction error: {str(e)}")
        return []

# Function to encode image to base64
def encode_image(file):
    try:
        file.seek(0)  # Ensure pointer is at start
        encoded_string = base64.b64encode(file.read()).decode("utf-8")
        return encoded_string
    except Exception as e:
        logger.error(f"Image encoding error: {str(e)}")
        raise ValueError(f"Failed to encode image: {str(e)}")

# Function to process file with MarkItDown
def process_with_markitdown(file):
    try:
        file.seek(0)  # Reset file pointer
        file_stream = io.BytesIO(file.read())  # Create binary file-like object
        result = md.convert_stream(file_stream)
        return result.text_content
    except ImportError as e:
        logger.error(f"MarkItDown dependency error: {str(e)}")
        return f"Missing dependencies for processing this file. Try: pip install 'markitdown[{file.name.split('.')[-1].lower()}]'"
    except ValueError as e:
        logger.error(f"MarkItDown processing error: {str(e)}")
        return f"Unsupported file format or invalid file: {str(e)}"
    except Exception as e:
        logger.error(f"Unexpected MarkItDown error: {str(e)}")
        return f"Error processing file with MarkItDown: {str(e)}"

# Function to process image or file with Grok
def extract_data(is_image=False, image_file=None):
    if not is_image or image_file is None:
        return None
    base64_image = encode_image(image_file)
    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{base64_image}",
                        "detail": "high",
                    },
                },
                {
                    "type": "text",
                    "text": "Please perform OCR and return all detected text in this image, as plain text."
                },
            ],
        },
    ]
    try:
        completion = client.chat.completions.create(
            model="grok-2-vision-latest",
            messages=messages,
            temperature=0.01,
        )
        logger.info(f"Grok OCR result: {completion.choices[0].message.content}")
        return completion.choices[0].message.content
    except Exception as e:
        logger.error(f"Grok extraction error: {str(e)}")
        return f"Error extracting data with Grok: {str(e)}. Check API key and network connection."

# Streamlit App
st.title("File and Image Data Extractor")
st.write("Upload Microsoft Office files (DOCX, XLSX, PPTX), PDFs, or images (JPEG/PNG) to extract structured data using Grok and MarkItDown.")

# Display optional dependencies info
st.info("Ensure all file formats are supported by installing dependencies: pip install 'markitdown[all]'")

# File uploader
uploaded_files = st.file_uploader(
    "Upload Files or Images",
    accept_multiple_files=True,
    type=["pdf", "docx", "xlsx", "pptx", "jpg", "jpeg", "png"]
)

# Process files
if uploaded_files:
    st.write("### Extracted Data")
    
    for file in uploaded_files:
        st.write(f"**Processing: {file.name}**")
        file_extension = file.name.split(".")[-1].lower()
        
        # Check if file is an image
        is_image = file_extension in ["jpg", "jpeg", "png"]
        
        parsed_text = None
        if is_image:
            st.info("Image detected. Using Grok Vision model for extraction.")
            result = extract_data(is_image=True, image_file=file)
            parsed_text = result  # Show Grok result in expander for images
        elif file_extension == "docx":
            content = process_with_markitdown(file)
            parsed_text = content
            if "Error" in content or "Missing dependencies" in content:
                st.error(content)
                continue
            images = extract_images_from_docx(file)
            image_results = []
            for idx, img in enumerate(images):
                st.info(f"Processing image {idx+1} in DOCX with Grok...")
                grok_result = extract_data(is_image=True, image_file=img)
                image_results.append((idx+1, grok_result))
            if image_results:
                combined = content + "\n\n" + "\n\n".join([
                    f"**Extracted from image {idx}:**\n{res}" for idx, res in image_results
                ])
                result = combined
            else:
                result = content
        elif file_extension == "pptx":
            content = process_with_markitdown(file)
            parsed_text = content
            if "Error" in content or "Missing dependencies" in content:
                st.error(content)
                continue
            images = extract_images_from_pptx(file)
            image_results = []
            for idx, img in enumerate(images):
                st.info(f"Processing image {idx+1} in PPTX with Grok...")
                grok_result = extract_data(is_image=True, image_file=img)
                image_results.append((idx+1, grok_result))
            if image_results:
                combined = content + "\n\n" + "\n\n".join([
                    f"**Extracted from image {idx}:**\n{res}" for idx, res in image_results
                ])
                result = combined
            else:
                result = content
        elif file_extension == "xlsx":
            content = process_with_markitdown(file)
            parsed_text = content
            if "Error" in content or "Missing dependencies" in content:
                st.error(content)
                continue
            images = extract_images_from_xlsx(file)
            image_results = []
            for idx, img in enumerate(images):
                st.info(f"Processing image {idx+1} in XLSX with Grok...")
                grok_result = extract_data(is_image=True, image_file=img)
                image_results.append((idx+1, grok_result))
            if image_results:
                combined = content + "\n\n" + "\n\n".join([
                    f"**Extracted from image {idx}:**\n{res}" for idx, res in image_results
                ])
                result = combined
            else:
                result = content
        elif file_extension == "pdf":
            content = process_with_markitdown(file)
            parsed_text = content
            if "Error" in content or "Missing dependencies" in content:
                st.error(content)
                continue
            images = extract_images_from_pdf(file)
            image_results = []
            for idx, img in enumerate(images):
                st.info(f"Processing image {idx+1} in PDF with Grok...")
                grok_result = extract_data(is_image=True, image_file=img)
                image_results.append((idx+1, grok_result))
            if image_results:
                combined = content + "\n\n" + "\n\n".join([
                    f"**Extracted from image {idx}:**\n{res}" for idx, res in image_results
                ])
                result = combined
            else:
                result = content
        else:
            content = process_with_markitdown(file)
            parsed_text = content
            if "Error" in content or "Missing dependencies" in content:
                st.error(content)
                continue
            st.info("Processed with MarkItDown. No Grok call for non-image files.")
            result = None  # Do not display parsed_text outside the expander
        
        # Show parsed text or Grok result in a closed expander (always)
        if is_image:
            if parsed_text and isinstance(parsed_text, str) and parsed_text.strip():
                with st.expander("Show Grok Extracted Info (Image)", expanded=False):
                    st.write(parsed_text)
            else:
                st.warning("No text detected in the image by Grok OCR.")
        elif parsed_text:
            with st.expander("Show parsed text (MarkItDown)", expanded=False):
                st.write(parsed_text)

        # Show Grok (image) extracted info in its own expander if available
        if 'image_results' in locals() and image_results:
            with st.expander("Show Grok Extracted Info (Images)", expanded=False):
                for idx, grok_result in image_results:
                    st.markdown(f"**Extracted from image {idx}:**")
                    st.markdown(grok_result)

        # Show Combined View expander if there is Grok content
        if 'image_results' in locals() and image_results:
            with st.expander("Combined View (Text + Grok Images)", expanded=False):
                # Insert text and Grok results in order: text, then each image result
                st.markdown("### Parsed Text (MarkItDown)")
                st.write(parsed_text)
                st.markdown("### Grok Extracted Info (Images)")
                for idx, grok_result in image_results:
                    st.markdown(f"**Extracted from image {idx}:**")
                    st.markdown(grok_result)